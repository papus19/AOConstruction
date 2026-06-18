"""
Extracteur de texte universel — supporte PDF, Word, Excel, PowerPoint, images, texte brut.
Fallback OCR automatique pour les PDFs scannés (images).
Utilisé par analyse.py et generateur_offres.py.
"""
import io
from typing import List
import streamlit as st


# ── Types acceptés par les file_uploader ──────────────────────────────
TYPES_ACCEPTES = None  # None = tous les types dans Streamlit

# Label et help uniformes pour tous les uploaders multi-documents
LABEL_UPLOAD = "Glissez vos documents ici ou cliquez pour sélectionner"
HELP_UPLOAD  = (
    "Tous types acceptés : PDF, Word (.docx/.doc), Excel (.xlsx/.xls), "
    "PowerPoint (.pptx), images (.png/.jpg/.tiff), texte (.txt/.csv) — "
    "Plusieurs fichiers simultanément"
)


# ════════════════════════════════════════════════════════════════════
# EXTRACTEURS PAR FORMAT
# ════════════════════════════════════════════════════════════════════

def _ocr_pdf(file_bytes: bytes) -> str:
    """
    OCR sur PDF scanné via pdf2image + pytesseract.
    Appelé automatiquement quand pypdf ne trouve pas de texte.
    
    Prérequis local :
      pip install pdf2image pytesseract Pillow
      + Tesseract installé sur le système
        Windows : https://github.com/UB-Mannheim/tesseract/wiki
        Mac     : brew install tesseract tesseract-lang
        Linux   : sudo apt install tesseract-ocr tesseract-ocr-fra
    """
    try:
        from pdf2image import convert_from_bytes
        import pytesseract

        # Sur Windows, si tesseract n'est pas dans le PATH, décommenter et ajuster :
        # pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

        images = convert_from_bytes(file_bytes, dpi=200)
        textes = []

        for i, img in enumerate(images, 1):
            texte_page = pytesseract.image_to_string(img, lang="fra+eng")
            if texte_page.strip():
                textes.append(f"=== Page {i} (OCR) ===\n{texte_page.strip()}")

        return "\n\n".join(textes) if textes else ""

    except ImportError:
        st.warning(
            "⚠️ OCR non disponible. "
            "Installez les dépendances : `pip install pdf2image pytesseract Pillow` "
            "puis Tesseract sur votre système."
        )
        return ""
    except Exception as e:
        st.warning(f"⚠️ OCR échoué : {e}")
        return ""


def _extraire_pdf(file_bytes: bytes) -> str:
    """
    Extrait le texte d'un PDF.
    - Essaie d'abord pypdf (rapide, pour les PDFs avec couche texte)
    - Si le résultat est vide ou trop court → PDF scanné → fallback OCR automatique
    """
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages  = [page.extract_text() or "" for page in reader.pages]
    texte  = "\n".join(pages).strip()

    # Seuil : moins de 50 caractères = probablement un PDF scanné
    if not texte or len(texte) < 50:
        st.info("📷 PDF scanné détecté — lancement de l'OCR…")
        texte = _ocr_pdf(file_bytes)

    return texte


def _extraire_docx(file_bytes: bytes) -> str:
    """Extrait le texte d'un fichier Word (.docx / .doc)"""
    from docx import Document

    doc   = Document(io.BytesIO(file_bytes))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]

    # Inclure les tableaux
    for table in doc.tables:
        for row in table.rows:
            ligne = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if ligne:
                paras.append(ligne)

    return "\n".join(paras)


def _extraire_xlsx(file_bytes: bytes) -> str:
    """Extrait le texte d'un fichier Excel (.xlsx / .xls)"""
    import openpyxl

    wb    = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    lignes = []

    for sheet in wb.worksheets:
        lignes.append(f"=== Feuille : {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            vals = [str(v) for v in row if v is not None]
            if vals:
                lignes.append(" | ".join(vals))

    return "\n".join(lignes)


def _extraire_pptx(file_bytes: bytes) -> str:
    """Extrait le texte d'une présentation PowerPoint (.pptx)"""
    from pptx import Presentation

    prs   = Presentation(io.BytesIO(file_bytes))
    texts = []

    for i, slide in enumerate(prs.slides, 1):
        texts.append(f"=== Diapositive {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

    return "\n".join(texts)


def _extraire_image(file_bytes: bytes) -> str:
    """
    Extrait le texte d'une image via OCR (pytesseract).
    Supporte : PNG, JPG, JPEG, TIFF, BMP, WEBP
    """
    try:
        import pytesseract
        from PIL import Image

        # Sur Windows, si tesseract n'est pas dans le PATH, décommenter et ajuster :
        # pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

        img  = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(img, lang="fra+eng")
        return text.strip()

    except ImportError:
        st.warning("⚠️ pytesseract ou Pillow non installé. `pip install pytesseract Pillow`")
        return "[Image — OCR non disponible]"
    except Exception as e:
        return f"[Image — extraction OCR échouée : {e}]"


def _extraire_txt(file_bytes: bytes) -> str:
    """Extrait le texte d'un fichier texte brut (.txt, .csv, .md, .rtf)"""
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


# ════════════════════════════════════════════════════════════════════
# FONCTION PRINCIPALE — UN FICHIER
# ════════════════════════════════════════════════════════════════════

def extraire_texte(uploaded_file) -> str:
    """
    Extrait le texte brut d'un fichier Streamlit UploadedFile.

    Formats supportés :
      - PDF      : pypdf + fallback OCR automatique si PDF scanné
      - Word     : .docx, .doc
      - Excel    : .xlsx, .xls
      - PowerPoint : .pptx
      - Images   : .png, .jpg, .jpeg, .tiff, .bmp, .webp  (OCR)
      - Texte    : .txt, .csv, .rtf, .md

    Retourne une chaîne vide avec un message d'avertissement si le format
    n'est pas reconnu ou si l'extraction échoue.
    """
    nom  = uploaded_file.name.lower()
    data = uploaded_file.read()
    uploaded_file.seek(0)  # Rembobiner pour usage ultérieur éventuel

    ext = nom.rsplit(".", 1)[-1] if "." in nom else ""

    try:
        if ext == "pdf":
            return _extraire_pdf(data)

        elif ext in ("docx", "doc"):
            return _extraire_docx(data)

        elif ext in ("xlsx", "xls"):
            return _extraire_xlsx(data)

        elif ext == "pptx":
            return _extraire_pptx(data)

        elif ext in ("png", "jpg", "jpeg", "tiff", "tif", "bmp", "webp"):
            return _extraire_image(data)

        elif ext in ("txt", "csv", "rtf", "md"):
            return _extraire_txt(data)

        else:
            # Tentative texte brut en dernier recours
            try:
                texte = _extraire_txt(data)
                if texte.strip():
                    return texte
            except Exception:
                pass
            st.warning(f"⚠️ Format .{ext} non reconnu — extraction impossible pour {uploaded_file.name}")
            return ""

    except Exception as e:
        st.warning(f"⚠️ Extraction partielle pour {uploaded_file.name} : {e}")
        return ""


# ════════════════════════════════════════════════════════════════════
# FONCTION PRINCIPALE — PLUSIEURS FICHIERS
# ════════════════════════════════════════════════════════════════════

def extraire_texte_multiple(uploaded_files) -> str:
    """
    Extrait et concatène le texte de plusieurs fichiers.
    Chaque fichier est séparé par un séparateur lisible.
    Le résultat est tronqué à 12 000 caractères pour les LLMs.

    Retourne le texte combiné, ou une chaîne vide si rien n'a été extrait.
    """
    if not uploaded_files:
        return ""

    parties = []

    for f in uploaded_files:
        texte = extraire_texte(f)
        if texte.strip():
            taille = f"{'%.1f' % (f.size / 1024)} KB"
            parties.append(
                f"━━━ DOCUMENT : {f.name} ({taille}) ━━━\n{texte}"
            )
        else:
            st.warning(f"⚠️ Aucun texte extrait de : {f.name}")

    if not parties:
        return ""

    combined = "\n\n".join(parties)

    # Tronquer intelligemment pour ne pas exploser le contexte LLM
    if len(combined) > 12000:
        combined = combined[:12000] + "\n\n[... Contenu tronqué à 12 000 caractères]"

    return combined


# ════════════════════════════════════════════════════════════════════
# FEEDBACK VISUEL
# ════════════════════════════════════════════════════════════════════

def feedback_fichiers(uploaded_files) -> None:
    """Affiche un résumé visuel des fichiers sélectionnés."""
    if not uploaded_files:
        st.caption("_Aucun document sélectionné_")
        return

    total_kb   = sum(f.size for f in uploaded_files) / 1024
    size_label = f"{total_kb / 1024:.1f} MB" if total_kb >= 1024 else f"{total_kb:.0f} KB"

    st.success(
        f"✅ **{len(uploaded_files)} fichier(s) prêt(s)** — "
        f"{', '.join(f.name for f in uploaded_files)} — {size_label} au total"
    )